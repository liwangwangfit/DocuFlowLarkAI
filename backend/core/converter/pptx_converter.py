"""
PPTX (PowerPoint) 转换器 - 将 PPTX 转换为 DOCX 以便飞书导入
"""
import tempfile
from pathlib import Path
from loguru import logger

from .base import BaseConverter, ConversionResult, DocumentType


class PPTXConverter(BaseConverter):
    """将 PowerPoint 文件转换为 DOCX 格式"""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type in [DocumentType.PPT, DocumentType.PPTX]

    async def convert(self, file_path: str) -> ConversionResult:
        try:
            return await self._convert_pptx_to_docx(file_path)
        except Exception as e:
            logger.error(f"PPTX转换失败: {e}")
            return ConversionResult(success=False, error_message=str(e))

    async def _convert_pptx_to_docx(self, file_path: str) -> ConversionResult:
        """提取 PPTX 幻灯片内容并写入 DOCX 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return ConversionResult(
                success=False,
                error_message="python-pptx 未安装，无法转换 PowerPoint 文件"
            )

        try:
            from docx import Document as DocxDocument
        except ImportError:
            return ConversionResult(
                success=False,
                error_message="python-docx 未安装，无法生成 DOCX 输出"
            )

        source = Path(file_path)
        prs = Presentation(file_path)
        doc = DocxDocument()
        doc.add_heading(source.stem, level=1)

        for slide_idx, slide in enumerate(prs.slides, 1):
            doc.add_heading(f"第 {slide_idx} 页", level=2)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            doc.add_paragraph(text)

                if shape.has_table:
                    table = shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)
                    doc_table = doc.add_table(rows=rows, cols=cols, style='Table Grid')
                    for r_idx, row in enumerate(table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            doc_table.rows[r_idx].cells[c_idx].text = cell.text

        output_dir = tempfile.mkdtemp(prefix="docuflow_pptx_")
        output_path = Path(output_dir) / f"{source.stem}.docx"
        doc.save(str(output_path))

        logger.info(f"PPTX转换完成: {source.name} -> {output_path.name}")
        return ConversionResult(
            success=True,
            output_path=str(output_path),
            metadata={
                "source_type": source.suffix.lower().lstrip("."),
                "title": source.stem,
                "slides": len(prs.slides),
            }
        )

    async def convert_to_markdown(self, file_path: str) -> ConversionResult:
        """将 PPTX 转换为 Markdown 字符串（备用路径）"""
        try:
            from pptx import Presentation
        except ImportError:
            return ConversionResult(
                success=False,
                error_message="python-pptx 未安装"
            )

        source = Path(file_path)
        prs = Presentation(file_path)
        lines = [f"# {source.stem}", ""]

        for slide_idx, slide in enumerate(prs.slides, 1):
            lines.append(f"## 第 {slide_idx} 页")
            lines.append("")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    header = [cell.text.strip() for cell in table.rows[0].cells]
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in list(table.rows)[1:]:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append("| " + " | ".join(cells) + " |")
            lines.append("")

        return ConversionResult(
            success=True,
            content="\n".join(lines),
            metadata={
                "source_type": source.suffix.lower().lstrip("."),
                "title": source.stem,
                "slides": len(prs.slides),
            }
        )
