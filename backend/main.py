"""
FastAPI 主入口
"""
import asyncio
import json
import time
from contextlib import asynccontextmanager
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple, Set

from fastapi import FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect, HTTPException, BackgroundTasks
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from config import get_config, BASE_DIR
from utils.logger import get_logger
from utils.system_resource import calculate_dynamic_workers

logger = get_logger(__name__)

# 全局状态
app_state = {
    "tasks": {},
    "connected_clients": set(),
    "system_status": "ready",  # ready, running, paused, error
    "stats": {
        "processed": 0,
        "failed": 0,
        "duplicate": 0,
        "tokens": 0,
        "api_calls": 0
    }
}

# 任务运行时状态锁（并发文件处理时保护共享状态）
task_state_lock = asyncio.Lock()


#  lifespan 上下文管理器
@asynccontextmanager
async def lifespan(app: FastAPI):
    """应用生命周期管理"""
    logger.info("=" * 50)
    logger.info("企业知识库迁移系统启动中...")
    logger.info("=" * 50)
    
    # 初始化数据库
    from models.database import init_db
    await init_db()
    logger.info("数据库初始化完成")
    
    # 初始化模板管理器
    from models.template import template_manager
    await template_manager.load_templates()
    logger.info("模板管理器初始化完成")
    
    # 测试飞书连接（检查是否有 user_access_token）
    try:
        from core.feishu.auth import feishu_oauth
        if feishu_oauth.is_authorized():
            logger.info("飞书 OAuth 授权有效")
        else:
            logger.warning("飞书 OAuth 未授权，请在配置面板完成授权")
    except Exception as e:
        logger.warning(f"飞书连接测试失败: {e}")
    
    yield
    
    # 关闭时清理
    logger.info("系统关闭中...")


app = FastAPI(
    title="企业知识库迁移系统",
    description="将本地文档迁移至飞书知识库",
    version="1.0.0",
    lifespan=lifespan
)

# CORS中间件
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ============ Pydantic模型 ============

class TaskCreateRequest(BaseModel):
    """创建任务请求"""
    name: str
    template_id: Optional[str] = None
    target_space_id: Optional[str] = None


class TaskResponse(BaseModel):
    """任务响应"""
    id: str
    name: str
    status: str
    progress: int
    created_at: str
    updated_at: Optional[str] = None


class ConfigUpdateRequest(BaseModel):
    """配置更新请求"""
    feishu: Optional[Dict[str, Any]] = None
    llm: Optional[Dict[str, Any]] = None
    mineru: Optional[Dict[str, Any]] = None


class LLMProcessRequest(BaseModel):
    """LLM处理请求"""
    content: str
    action: str  # clean, summarize, classify, quality_check
    template_id: Optional[str] = None


class FeishuUploadRequest(BaseModel):
    """飞书上传请求"""
    document_id: str
    title: str
    content: str
    parent_node_token: Optional[str] = None


# ============ WebSocket管理 ============

class ConnectionManager:
    """WebSocket连接管理器"""
    
    def __init__(self):
        self.active_connections: List[WebSocket] = []
    
    async def connect(self, websocket: WebSocket):
        await websocket.accept()
        self.active_connections.append(websocket)
        logger.info(f"WebSocket客户端连接，当前连接数: {len(self.active_connections)}")
    
    def disconnect(self, websocket: WebSocket):
        if websocket in self.active_connections:
            self.active_connections.remove(websocket)
        logger.info(f"WebSocket客户端断开，当前连接数: {len(self.active_connections)}")
    
    async def broadcast(self, message: dict):
        """广播消息给所有客户端"""
        disconnected = []
        for connection in self.active_connections:
            try:
                await connection.send_json(message)
            except Exception:
                disconnected.append(connection)
        
        # 清理断开的连接
        for conn in disconnected:
            self.disconnect(conn)
    
    async def send_log(self, source: str, message: str, level: str = "info"):
        """发送日志消息"""
        await self.broadcast({
            "type": "log",
            "data": {
                "time": datetime.now().strftime("%H:%M:%S"),
                "source": source,
                "message": message,
                "level": level
            }
        })
    
    async def send_progress(self, task_id: str, progress: int, phase: str):
        """发送进度更新"""
        await self.broadcast({
            "type": "progress",
            "data": {
                "task_id": task_id,
                "progress": progress,
                "phase": phase
            }
        })
    
    async def send_stats(self, stats: dict):
        """发送统计更新"""
        # 确保字段名与前端一致
        data = {
            "processed": stats.get("processed", 0),
            "failed": stats.get("failed", 0),
            "duplicate": stats.get("duplicate", 0),
            "tokens": stats.get("tokens", 0),
            "api": stats.get("api_calls", 0)  # 前端使用 api 字段名
        }
        await self.broadcast({
            "type": "stats",
            "data": data
        })
    
    async def send_panorama(self, task_id: str, space_structure: dict, file_status: list):
        """
        发送全景图数据
        
        space_structure: 知识库空间结构
        file_status: 文件处理状态列表
        """
        await self.broadcast({
            "type": "panorama",
            "data": {
                "task_id": task_id,
                "space_structure": space_structure,
                "file_status": file_status,
                "timestamp": datetime.now().isoformat()
            }
        })
    
    async def send_chart_data(self, task_id: str, chart_type: str, data: list):
        """
        发送图表数据
        
        chart_type: 'line' | 'bar' | 'pie'
        data: 图表数据点列表
        """
        await self.broadcast({
            "type": "chart",
            "data": {
                "task_id": task_id,
                "chart_type": chart_type,
                "data": data,
                "timestamp": datetime.now().isoformat()
            }
        })


manager = ConnectionManager()


# ============ API路由 ============

@app.get("/")
async def root():
    """返回主页面"""
    html_path = BASE_DIR / "frontend" / "index.html"
    if html_path.exists():
        return FileResponse(html_path)
    return {"message": "企业知识库迁移系统 API", "version": "1.0.0"}


@app.get("/api/status")
async def get_status():
    """获取系统状态"""
    return {
        "status": app_state["system_status"],
        "stats": app_state["stats"],
        "active_tasks": len(app_state["tasks"]),
        "connected_clients": len(manager.active_connections)
    }


@app.get("/api/config")
async def get_configuration():
    """获取配置信息"""
    config = get_config()
    
    # 获取当前启用的 provider 和 API key 状态（含每个 provider 的 key 映射）
    active_provider = None
    has_api_key = False
    has_api_key_map: Dict[str, bool] = {}
    for name, provider in config.llm.providers.items():
        # provider 可能是 dict 或 Pydantic 模型
        if isinstance(provider, dict):
            provider_has_key = bool(provider.get("api_key"))
            has_api_key_map[name] = provider_has_key
            if active_provider is None and provider.get("enabled"):
                active_provider = name
                has_api_key = provider_has_key
        else:
            # Pydantic 模型
            provider_has_key = bool(getattr(provider, "api_key", None))
            has_api_key_map[name] = provider_has_key
            if active_provider is None and getattr(provider, "enabled", False):
                active_provider = name
                has_api_key = provider_has_key
    
    return {
        "feishu": {
            "app_id": config.feishu.app_id[:10] + "***" if config.feishu.app_id else None,
            "has_app_id": bool(config.feishu.app_id),
            "has_app_secret": bool(config.feishu.app_secret),
            "base_url": config.feishu.base_url,
            "mcp_url": config.feishu.mcp_url
        },
        "llm": {
            "provider": active_provider,
            "has_api_key": has_api_key,
            "has_api_key_map": has_api_key_map,
            "content_clean": config.llm.content_clean,
            "quality_check": config.llm.quality_check,
            "min_score": config.llm.min_score
        },
        "mineru": {
            "use_local": config.mineru.use_local,
            "local_url": config.mineru.local_url
        }
    }


class FeishuCredentialTestRequest(BaseModel):
    app_id: Optional[str] = None
    app_secret: Optional[str] = None


@app.post("/api/feishu/credentials/test")
async def test_feishu_credentials(request: FeishuCredentialTestRequest):
    """测试飞书 App ID / App Secret 是否有效（不依赖用户OAuth授权）"""
    import httpx

    config = get_config()
    app_id = (request.app_id or config.feishu.app_id or "").strip()
    app_secret = (request.app_secret or config.feishu.app_secret or "").strip()
    if not app_id or not app_secret:
        raise HTTPException(status_code=400, detail="请先填写 App ID 和 App Secret")

    token_url = f"{config.feishu.base_url}/auth/v3/tenant_access_token/internal"
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            resp = await client.post(
                token_url,
                headers={"Content-Type": "application/json; charset=utf-8"},
                json={"app_id": app_id, "app_secret": app_secret},
            )
            data = resp.json()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"请求飞书鉴权接口失败: {e}")

    if data.get("code") != 0:
        msg = data.get("msg") or data.get("error_description") or str(data)
        raise HTTPException(status_code=400, detail=f"飞书应用凭据无效: {msg}")

    return {"message": "飞书应用连接测试成功"}


@app.post("/api/config")
async def update_configuration(request: ConfigUpdateRequest):
    """更新配置 - 直接写入YAML文件，加载时直接读取"""
    import yaml
    from config import CONFIG_DIR, get_config, protect_config_data, reload_config
    
    current_config = get_config()
    
    def _model_to_dict(data: Any) -> Dict[str, Any]:
        if data is None:
            return {}
        if isinstance(data, dict):
            return dict(data)
        if hasattr(data, "model_dump"):
            return data.model_dump()
        if hasattr(data, "dict"):
            return data.dict()
        return {}
    
    if request.feishu:
        # 合并已有配置，避免前端未展示字段被覆盖
        feishu_data = _model_to_dict(current_config.feishu)
        feishu_data.update(dict(request.feishu))
        feishu_path = CONFIG_DIR / "feishu.yaml"
        with open(feishu_path, 'w', encoding='utf-8') as f:
            yaml.safe_dump(
                protect_config_data(feishu_data),
                f,
                allow_unicode=True,
                default_flow_style=False,
                sort_keys=False,
            )
    
    if request.llm:
        # 定义默认provider配置
        default_providers = {
            "deepseek": {"enabled": False, "model": "deepseek-chat", "base_url": "https://api.deepseek.com", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "kimi": {"enabled": False, "model": "moonshot-v1-8k", "base_url": "https://api.moonshot.cn", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "openai": {"enabled": False, "model": "gpt-4-turbo-preview", "base_url": "https://api.openai.com", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "claude": {"enabled": False, "model": "claude-3-sonnet-20240229", "base_url": "https://api.anthropic.com", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "qwen": {"enabled": False, "model": "qwen-max", "base_url": "https://dashscope.aliyuncs.com/api/v1", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "ollama": {"enabled": False, "model": "llama2-chinese", "base_url": "http://localhost:11434", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""},
            "mock": {"enabled": False, "model": "mock", "base_url": "http://localhost/mock", "temperature": 0.3, "max_tokens": 4000, "timeout": 60, "api_key": ""}
        }
        
        # 合并已有LLM配置，避免隐藏字段（如 min_score）被覆盖
        llm_data = _model_to_dict(current_config.llm)
        if "providers" in llm_data:
            llm_data.pop("providers")
        llm_data.update(dict(request.llm))
        
        # 从前端数据中提取信息
        provider = llm_data.pop("provider", None)
        api_key = llm_data.pop("api_key", None)
        current_providers = getattr(current_config.llm, "providers", {}) or {}
        
        # 构建 providers 结构 - 使用纯Python字典确保YAML序列化正确
        providers = {}
        
        # 设置当前选中的provider
        if provider and provider in default_providers:
            providers[provider] = default_providers[provider].copy()
            providers[provider]["enabled"] = True
            if api_key:
                providers[provider]["api_key"] = api_key
            else:
                # 未输入新 key 时，保留已有 key（避免切换/保存时丢失）
                existing_provider = current_providers.get(provider)
                if isinstance(existing_provider, dict):
                    existing_key = existing_provider.get("api_key")
                else:
                    existing_key = getattr(existing_provider, "api_key", None) if existing_provider else None
                if existing_key:
                    providers[provider]["api_key"] = existing_key
        
        # 其他 provider 保持未启用状态（从默认配置复制）
        for name, cfg in default_providers.items():
            if name != provider:
                providers[name] = cfg.copy()
                providers[name]["enabled"] = False
                existing_provider = current_providers.get(name)
                if isinstance(existing_provider, dict):
                    existing_key = existing_provider.get("api_key")
                else:
                    existing_key = getattr(existing_provider, "api_key", None) if existing_provider else None
                if existing_key:
                    providers[name]["api_key"] = existing_key
        
        llm_data["providers"] = providers
        
        # 保存到YAML文件
        llm_path = CONFIG_DIR / "llm.yaml"
        with open(llm_path, 'w', encoding='utf-8') as f:
            yaml.safe_dump(
                protect_config_data(llm_data),
                f,
                allow_unicode=True,
                default_flow_style=False,
                sort_keys=False,
            )
        
        logger.info(f"LLM配置已保存到 {llm_path}: provider={provider}, has_api_key={bool(api_key)}")
    
    if request.mineru:
        mineru_path = CONFIG_DIR / "mineru.yaml"
        with open(mineru_path, 'w', encoding='utf-8') as f:
            yaml.safe_dump(
                protect_config_data(dict(request.mineru)),
                f,
                allow_unicode=True,
                default_flow_style=False,
                sort_keys=False,
            )
    
    # 重新加载配置
    reload_config()
    
    return {"message": "配置已更新"}


@app.post("/api/config/clear")
async def clear_configuration():
    """清空敏感配置并登出授权信息"""
    import yaml
    from config import CONFIG_DIR, get_config, protect_config_data, reload_config
    from core.feishu.auth import feishu_oauth

    current_config = get_config()

    def _model_to_dict(data: Any) -> Dict[str, Any]:
        if data is None:
            return {}
        if isinstance(data, dict):
            return dict(data)
        if hasattr(data, "model_dump"):
            return data.model_dump()
        if hasattr(data, "dict"):
            return data.dict()
        return {}

    # 飞书敏感配置清空
    feishu_data = _model_to_dict(current_config.feishu)
    feishu_data["app_id"] = ""
    feishu_data["app_secret"] = ""
    feishu_data["encrypt_key"] = None
    feishu_data["verification_token"] = None
    with open(CONFIG_DIR / "feishu.yaml", "w", encoding="utf-8") as f:
        yaml.safe_dump(
            protect_config_data(feishu_data),
            f,
            allow_unicode=True,
            default_flow_style=False,
            sort_keys=False,
        )

    # LLM 全provider关闭并清空 API Key
    llm_data = _model_to_dict(current_config.llm)
    providers = llm_data.get("providers", {}) or {}
    sanitized_providers: Dict[str, Dict[str, Any]] = {}
    for name, provider in providers.items():
        if isinstance(provider, dict):
            provider_data = dict(provider)
        elif hasattr(provider, "model_dump"):
            provider_data = provider.model_dump()
        elif hasattr(provider, "dict"):
            provider_data = provider.dict()
        else:
            provider_data = {}
        provider_data["enabled"] = False
        provider_data["api_key"] = ""
        sanitized_providers[name] = provider_data
    llm_data["providers"] = sanitized_providers
    with open(CONFIG_DIR / "llm.yaml", "w", encoding="utf-8") as f:
        yaml.safe_dump(
            protect_config_data(llm_data),
            f,
            allow_unicode=True,
            default_flow_style=False,
            sort_keys=False,
        )

    # MinerU 云端 Key 清空
    mineru_data = _model_to_dict(current_config.mineru)
    mineru_data["cloud_api_key"] = ""
    with open(CONFIG_DIR / "mineru.yaml", "w", encoding="utf-8") as f:
        yaml.safe_dump(
            protect_config_data(mineru_data),
            f,
            allow_unicode=True,
            default_flow_style=False,
            sort_keys=False,
        )

    # 清空 OAuth token
    feishu_oauth.clear_token()
    reload_config()
    return {"message": "配置与授权信息已清空"}


@app.post("/api/tasks")
async def create_task(request: TaskCreateRequest):
    """创建迁移任务"""
    import uuid
    
    task_id = str(uuid.uuid4())[:8]
    task = {
        "id": task_id,
        "name": request.name,
        "status": "pending",
        "progress": 0,
        "template_id": request.template_id,
        "target_space_id": request.target_space_id,
        "created_at": datetime.now().isoformat(),
        "updated_at": None,
        "files": [],
        "duplicates": [],
        "results": []
    }
    
    app_state["tasks"][task_id] = task
    
    await manager.send_log("Task", f"创建任务: {request.name} (ID: {task_id})", "info")
    
    return task


@app.get("/api/tasks")
async def list_tasks():
    """获取任务列表"""
    return list(app_state["tasks"].values())


@app.get("/api/tasks/{task_id}")
async def get_task(task_id: str):
    """获取任务详情"""
    if task_id not in app_state["tasks"]:
        raise HTTPException(status_code=404, detail="任务不存在")
    return app_state["tasks"][task_id]


def _normalize_name_for_compare(name: str) -> str:
    """文件名规范化（忽略扩展名、大小写）"""
    if not name:
        return ""
    raw = name.strip().casefold()
    stem = Path(name).stem.strip().casefold()
    return stem or raw


async def _resolve_duplicate_check_space_id(task: Dict[str, Any]) -> Optional[str]:
    """
    解析用于重名检查的知识空间ID。
    优先级：任务指定空间 > 配置默认空间 > 同名空间（模板名/任务名）。
    """
    if task.get("target_space_id"):
        return task["target_space_id"]

    try:
        config = get_config()
        default_space_id = getattr(config.feishu, "default_space_id", None)
        if default_space_id:
            return default_space_id
    except Exception as e:
        logger.warning(f"读取默认知识空间失败: {e}")

    candidate_name = (task.get("name") or "").strip()
    template_id = task.get("template_id")
    if template_id:
        try:
            from models.template import template_manager
            template = await template_manager.get_template(template_id)
            if template and template.get("name"):
                candidate_name = str(template["name"]).strip()
        except Exception as e:
            logger.warning(f"读取模板信息失败（重名检查降级）: {e}")

    if not candidate_name:
        return None

    try:
        from core.feishu.wiki_api import wiki_api

        spaces = await wiki_api.list_spaces(page_size=100)
        for space in spaces:
            if (space.get("name") or "").strip() == candidate_name:
                return space.get("space_id")
    except Exception as e:
        logger.warning(f"按名称查找知识空间失败（重名检查降级）: {e}")

    return None


async def _collect_existing_kb_name_sets(space_id: str) -> Tuple[Set[str], Set[str]]:
    """收集知识空间内已存在节点名（原名集合 + 规范化集合）"""
    from core.feishu.wiki_api import wiki_api

    existing_full: Set[str] = set()
    existing_normalized: Set[str] = set()
    nodes = await wiki_api.get_space_nodes_flat(space_id)
    for node in nodes:
        title = (node.get("title") or "").strip()
        if not title:
            continue
        existing_full.add(title.casefold())
        normalized = _normalize_name_for_compare(title)
        if normalized:
            existing_normalized.add(normalized)
    return existing_full, existing_normalized


@app.post("/api/tasks/{task_id}/upload")
async def upload_files(task_id: str, files: List[UploadFile] = File(...)):
    """上传文件到任务"""
    if task_id not in app_state["tasks"]:
        raise HTTPException(status_code=404, detail="任务不存在")
    
    task = app_state["tasks"][task_id]
    task.setdefault("duplicates", [])
    uploaded = []
    duplicates = []
    
    from config import CACHE_DIR
    task_cache_dir = CACHE_DIR / task_id
    task_cache_dir.mkdir(exist_ok=True)

    existing_task_full: Set[str] = set()
    existing_task_normalized: Set[str] = set()
    for info in task.get("files", []):
        name = info.get("name", "")
        if not name:
            continue
        existing_task_full.add(name.casefold())
        normalized = _normalize_name_for_compare(name)
        if normalized:
            existing_task_normalized.add(normalized)

    for info in task.get("duplicates", []):
        name = info.get("name", "")
        if not name:
            continue
        existing_task_full.add(name.casefold())
        normalized = _normalize_name_for_compare(name)
        if normalized:
            existing_task_normalized.add(normalized)

    kb_existing_full: Set[str] = set()
    kb_existing_normalized: Set[str] = set()
    duplicate_check_space_id = await _resolve_duplicate_check_space_id(task)
    if duplicate_check_space_id:
        try:
            kb_existing_full, kb_existing_normalized = await _collect_existing_kb_name_sets(duplicate_check_space_id)
            await manager.send_log(
                "Upload",
                f"重名检查：已加载知识空间 {duplicate_check_space_id}，现有节点 {len(kb_existing_full)} 个",
                "info",
            )
        except Exception as e:
            await manager.send_log("Upload", f"重名检查降级（无法读取知识空间）: {e}", "warn")
    
    for file in files:
        file_name = (file.filename or "").strip()
        if not file_name:
            file_name = f"unnamed_{int(time.time() * 1000)}"

        file_name_full = file_name.casefold()
        file_name_normalized = _normalize_name_for_compare(file_name)
        duplicated_in_kb = (
            file_name_full in kb_existing_full
            or (file_name_normalized and file_name_normalized in kb_existing_normalized)
        )
        duplicated_in_task = (
            file_name_full in existing_task_full
            or (file_name_normalized and file_name_normalized in existing_task_normalized)
        )

        if duplicated_in_kb or duplicated_in_task:
            duplicate_info = {
                "name": file_name,
                "reason": "existing_in_kb" if duplicated_in_kb else "duplicate_in_task",
            }
            duplicates.append(duplicate_info)

            if file_name_full not in existing_task_full:
                task["duplicates"].append(duplicate_info)
                existing_task_full.add(file_name_full)
                if file_name_normalized:
                    existing_task_normalized.add(file_name_normalized)

            reason_text = "知识库中已存在同名文件" if duplicated_in_kb else "任务中已存在同名文件"
            await manager.send_log("Upload", f"跳过重复文件: {file_name}（{reason_text}）", "warn")
            continue

        file_path = task_cache_dir / file_name
        content = await file.read()
        with open(file_path, 'wb') as f:
            f.write(content)
        
        file_info = {
            "name": file_name,
            "path": str(file_path),
            "size": len(content),
            "type": file.content_type
        }
        task["files"].append(file_info)
        uploaded.append(file_info)

        existing_task_full.add(file_name_full)
        if file_name_normalized:
            existing_task_normalized.add(file_name_normalized)
        
        await manager.send_log("Upload", f"上传文件: {file_name} ({len(content)} bytes)", "info")
    
    return {
        "uploaded": uploaded,
        "duplicates": duplicates,
        "total": len(task["files"]),
        "duplicate_total": len(task.get("duplicates", [])),
    }


@app.post("/api/tasks/{task_id}/start")
async def start_task(task_id: str, background_tasks: BackgroundTasks):
    """开始执行任务"""
    if task_id not in app_state["tasks"]:
        raise HTTPException(status_code=404, detail="任务不存在")
    
    task = app_state["tasks"][task_id]
    
    if task["status"] == "running":
        raise HTTPException(status_code=400, detail="任务正在运行中")
    
    task["status"] = "running"
    task["updated_at"] = datetime.now().isoformat()
    
    # 在后台执行任务
    background_tasks.add_task(run_migration_task, task_id)
    
    await manager.send_log("Task", f"开始执行任务: {task_id}", "info")
    
    return {"message": "任务已开始", "task_id": task_id}


@app.post("/api/tasks/{task_id}/cancel")
async def cancel_task(task_id: str):
    """取消任务"""
    if task_id not in app_state["tasks"]:
        raise HTTPException(status_code=404, detail="任务不存在")
    
    task = app_state["tasks"][task_id]
    task["status"] = "cancelled"
    task["updated_at"] = datetime.now().isoformat()
    
    await manager.send_log("Task", f"取消任务: {task_id}", "warn")
    
    return {"message": "任务已取消"}


@app.delete("/api/tasks/{task_id}")
async def delete_task(task_id: str):
    """删除任务"""
    if task_id not in app_state["tasks"]:
        raise HTTPException(status_code=404, detail="任务不存在")
    
    del app_state["tasks"][task_id]
    
    # 清理缓存
    from config import CACHE_DIR
    task_cache_dir = CACHE_DIR / task_id
    if task_cache_dir.exists():
        import shutil
        shutil.rmtree(task_cache_dir)
    
    return {"message": "任务已删除"}


@app.post("/api/llm/process")
async def llm_process(request: LLMProcessRequest):
    """LLM处理接口"""
    from core.llm.processor import LLMProcessor
    
    processor = LLMProcessor()
    
    try:
        if request.action == "clean":
            result = await processor.clean_content(request.content)
        elif request.action == "summarize":
            result = await processor.summarize(request.content)
        elif request.action == "classify":
            result = await processor.classify(request.content, request.template_id)
        elif request.action == "quality_check":
            result = await processor.quality_check(request.content)
        else:
            raise HTTPException(status_code=400, detail="未知的操作类型")
        
        # 更新统计
        app_state["stats"]["tokens"] += result.get("tokens", 0)
        await manager.send_stats(app_state["stats"])
        
        return result
    
    except Exception as e:
        logger.error(f"LLM处理失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/feishu/upload")
async def feishu_upload(request: FeishuUploadRequest):
    """飞书文档上传"""
    from core.feishu.document_api import FeishuDocumentAPI
    
    api = FeishuDocumentAPI()
    
    try:
        # 创建文档
        doc_id = await api.create_document(
            title=request.title,
            folder_token=request.parent_node_token
        )
        
        # 转换内容
        blocks = await api.convert_markdown_to_blocks(request.content)
        
        # 批量插入
        await api.batch_insert_blocks(doc_id, blocks)
        
        # 更新统计
        app_state["stats"]["api_calls"] += 3
        await manager.send_stats(app_state["stats"])
        
        return {
            "success": True,
            "document_id": doc_id,
            "title": request.title,
            "blocks_count": len(blocks)
        }
    
    except Exception as e:
        logger.error(f"飞书上传失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ============ 模板管理 API ============

@app.get("/api/templates")
async def list_templates():
    """获取所有模板列表"""
    from models.template import template_manager
    templates = await template_manager.list_templates()
    return templates


@app.get("/api/templates/{template_id}")
async def get_template(template_id: str):
    """获取模板详情"""
    from models.template import template_manager
    template = await template_manager.get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    return template


@app.post("/api/templates")
async def create_template(request: dict):
    """创建新模板"""
    from models.template import template_manager
    template = await template_manager.create_template(request)
    return template


@app.put("/api/templates/{template_id}")
async def update_template(template_id: str, request: dict):
    """更新模板基本信息"""
    from models.template import template_manager
    template = await template_manager.update_template(template_id, request)
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    return template


@app.delete("/api/templates/{template_id}")
async def delete_template(template_id: str):
    """删除模板"""
    from models.template import template_manager
    success = await template_manager.delete_template(template_id)
    if not success:
        raise HTTPException(status_code=404, detail="模板不存在")
    return {"message": "模板已删除"}


@app.post("/api/templates/import")
async def import_template(file: UploadFile = File(...)):
    """
    导入模板 JSON 文件
    
    上传 JSON 文件创建新模板
    """
    from models.template import template_manager
    
    if not file.filename.endswith('.json'):
        raise HTTPException(status_code=400, detail="只支持 JSON 文件")
    
    try:
        content = await file.read()
        data = json.loads(content.decode('utf-8'))
        
        # 验证必要字段
        if "name" not in data:
            raise HTTPException(status_code=400, detail="JSON 缺少 name 字段")
        
        # 生成新 ID，避免冲突
        if "id" in data:
            data["id"] = f"{data['id']}_import_{int(time.time())}"
        
        template = await template_manager.create_template(data)
        return {
            "message": "模板导入成功",
            "template": template
        }
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="无效的 JSON 文件")
    except Exception as e:
        logger.error(f"导入模板失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/templates/{template_id}/export")
async def export_template(template_id: str):
    """
    导出模板为 JSON 文件
    
    下载模板结构的 JSON 文件
    """
    from models.template import template_manager
    
    template = await template_manager.get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    
    # 创建导出内容（去掉内部字段）
    export_data = {
        "name": template["name"],
        "description": template.get("description", ""),
        "structure": template.get("structure", [])
    }
    
    # 返回 JSON 文件
    from fastapi.responses import StreamingResponse
    import io
    
    json_content = json.dumps(export_data, ensure_ascii=False, indent=2)
    json_bytes = json_content.encode('utf-8')
    
    return StreamingResponse(
        io.BytesIO(json_bytes),
        media_type="application/json",
        headers={
            "Content-Disposition": f"attachment; filename=template_{template_id}.json"
        }
    )


@app.post("/api/templates/{template_id}/export/save")
async def export_template_to_local(template_id: str):
    """
    导出模板到本地文件（桌面模式友好）

    在 data/exports 目录生成 JSON 文件并返回路径。
    """
    from models.template import template_manager

    template = await template_manager.get_template(template_id)
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")

    export_data = {
        "name": template["name"],
        "description": template.get("description", ""),
        "structure": template.get("structure", []),
    }

    export_dir = BASE_DIR / "data" / "exports"
    export_dir.mkdir(parents=True, exist_ok=True)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_name = f"template_{template_id}_{ts}.json"
    file_path = export_dir / file_name

    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(export_data, f, ensure_ascii=False, indent=2)

    return {
        "message": "模板导出成功",
        "file_name": file_name,
        "file_path": str(file_path),
    }


from pydantic import BaseModel
from typing import Optional

class NodeData(BaseModel):
    name: str
    type: str = "folder"

class AddNodeRequest(BaseModel):
    parent_id: Optional[str] = None
    node: NodeData

class UpdateNodeRequest(BaseModel):
    name: str
    type: str

@app.post("/api/templates/{template_id}/nodes")
async def add_template_node(template_id: str, request: AddNodeRequest):
    """添加模板节点"""
    from models.template import template_manager
    
    template = await template_manager.add_node(template_id, request.parent_id, request.node.dict())
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    return template


@app.put("/api/templates/{template_id}/nodes/{node_id}")
async def update_template_node(template_id: str, node_id: str, request: UpdateNodeRequest):
    """更新模板节点"""
    from models.template import template_manager
    template = await template_manager.update_node(template_id, node_id, request.dict())
    if not template:
        raise HTTPException(status_code=404, detail="模板或节点不存在")
    return template


@app.delete("/api/templates/{template_id}/nodes/{node_id}")
async def delete_template_node(template_id: str, node_id: str):
    """删除模板节点"""
    from models.template import template_manager
    template = await template_manager.delete_node(template_id, node_id)
    if not template:
        raise HTTPException(status_code=404, detail="模板不存在")
    return template


@app.get("/api/templates/{template_id}/format")
async def format_template_for_llm(template_id: str):
    """获取LLM格式的模板文本"""
    from models.template import template_manager
    text = template_manager.format_for_llm(template_id)
    return {"text": text}


@app.get("/api/feishu/spaces")
async def list_feishu_spaces():
    """获取飞书知识空间列表"""
    from core.feishu.wiki_api import FeishuWikiAPI
    
    api = FeishuWikiAPI()
    
    try:
        spaces = await api.list_spaces()
        return spaces
    except Exception as e:
        logger.error(f"获取知识空间列表失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/feishu/spaces")
async def create_feishu_space(name: str, description: str = ""):
    """创建飞书知识空间"""
    from core.feishu.wiki_api import FeishuWikiAPI
    
    api = FeishuWikiAPI()
    
    try:
        space_id = await api.create_space(name, description)
        return {"space_id": space_id, "name": name}
    except Exception as e:
        logger.error(f"创建知识空间失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ============ OAuth 授权 API ============

@app.get("/api/auth/status")
async def get_auth_status():
    """获取当前授权状态"""
    from core.feishu.auth import feishu_oauth
    return {
        "authorized": feishu_oauth.is_authorized(),
        "token_info": feishu_oauth.get_token_info()
    }


@app.get("/api/auth/url")
async def get_auth_url(redirect_uri: str = "http://127.0.0.1:8000/auth/callback"):
    """获取飞书 OAuth 授权 URL"""
    from core.feishu.auth import feishu_oauth
    from config import get_config

    feishu_cfg = get_config().feishu
    app_id = (feishu_cfg.app_id or "").strip()
    app_secret = (feishu_cfg.app_secret or "").strip()
    if not app_id or not app_secret:
        raise HTTPException(
            status_code=400,
            detail="请先在配置页面填写并保存飞书 App ID 和 App Secret，再进行授权"
        )

    # 同步最新配置，避免配置热更新后使用旧凭据
    feishu_oauth.config = feishu_cfg
    auth_url = feishu_oauth.get_authorize_url(redirect_uri=redirect_uri)
    return {
        "auth_url": auth_url,
        "redirect_uri": redirect_uri
    }


@app.post("/api/auth/exchange")
async def exchange_auth_code(code: str, redirect_uri: str = "http://127.0.0.1:8000/auth/callback"):
    """
    用授权码换取 user_access_token
    
    前端在完成授权后，从回调 URL 中获取 code，调用此接口
    """
    from core.feishu.auth import feishu_oauth
    logger.info(f"收到授权码交换请求: code={code[:10]}..., redirect_uri={redirect_uri}")
    try:
        token_data = await feishu_oauth.exchange_code_for_token(code, redirect_uri)
        logger.info(f"授权成功: expires_in={token_data.get('expires_in')}")
        return {
            "success": True,
            "message": "授权成功",
            "expires_in": token_data.get("expires_in")
        }
    except Exception as e:
        logger.error(f"授权失败: {e}")
        raise HTTPException(status_code=400, detail=f"授权失败: {str(e)}")


@app.post("/api/auth/refresh")
async def refresh_auth_token():
    """刷新 user_access_token"""
    from core.feishu.auth import feishu_oauth
    import time
    
    token_info = feishu_oauth.get_token_info()
    
    # 检查是否有 refresh_token
    if not token_info["has_refresh_token"]:
        logger.error("刷新失败: 没有 refresh_token")
        raise HTTPException(
            status_code=400, 
            detail="没有刷新令牌。请重新授权，并确保勾选'长期访问权限'选项"
        )
    
    # 检查 refresh_token 是否过期
    if not token_info["refresh_token_valid"]:
        logger.error("刷新失败: refresh_token 已过期")
        raise HTTPException(
            status_code=400, 
            detail="刷新令牌已过期（超过7天）。请重新授权"
        )
    
    try:
        token_data = await feishu_oauth.refresh_access_token()
        return {
            "success": True,
            "message": "刷新成功",
            "expires_in": token_data.get("expires_in"),
            "token_info": feishu_oauth.get_token_info()
        }
    except Exception as e:
        error_msg = str(e)
        logger.error(f"刷新失败: {error_msg}")
        
        # 提供更友好的错误信息
        if "20064" in error_msg or "20073" in error_msg:
            detail = "刷新令牌已被使用或撤销。请重新授权"
        elif "20037" in error_msg:
            detail = "刷新令牌已过期（超过7天）。请重新授权"
        else:
            detail = f"刷新失败: {error_msg}"
        
        raise HTTPException(status_code=400, detail=detail)


@app.post("/api/auth/logout")
async def logout():
    """清除授权（登出）"""
    from core.feishu.auth import feishu_oauth
    feishu_oauth.clear_token()
    return {"message": "已登出"}


# ============ WebSocket ============

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    """WebSocket端点"""
    await manager.connect(websocket)
    
    try:
        # 发送初始状态
        await websocket.send_json({
            "type": "status",
            "data": {
                "status": app_state["system_status"],
                "stats": app_state["stats"]
            }
        })
        
        while True:
            # 接收客户端消息
            data = await websocket.receive_json()
            
            if data.get("action") == "ping":
                await websocket.send_json({"type": "pong"})
            
            elif data.get("action") == "get_logs":
                # 返回最近日志
                pass
            
    except WebSocketDisconnect:
        manager.disconnect(websocket)
    except Exception as e:
        logger.error(f"WebSocket错误: {e}")
        manager.disconnect(websocket)


# ============ OAuth 回调页面 ============

@app.get("/auth/callback")
async def auth_callback_page(code: str = None, state: str = None, error: str = None):
    """
    OAuth 回调页面

    兼容两种模式：
    1. 弹窗模式：postMessage 通知父窗口后自动关闭
    2. 桌面/无弹窗模式：本页直接完成 token 交换后跳回首页
    """
    import html as _html
    from core.feishu.auth import feishu_oauth

    def _render_page(title: str, message: str, is_success: bool, script: str) -> HTMLResponse:
        color = "#16a34a" if is_success else "#dc2626"
        icon = "✅" if is_success else "❌"
        return HTMLResponse(
            f"""
            <!DOCTYPE html>
            <html>
            <head>
                <title>{_html.escape(title)}</title>
                <style>
                    body {{ font-family: system-ui, sans-serif; display: flex; justify-content: center; align-items: center; height: 100vh; margin: 0; background: #f5f5f5; }}
                    .card {{ background: white; padding: 2rem; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); text-align: center; max-width: 460px; }}
                    .msg {{ color: {color}; margin: 1rem 0; word-break: break-word; }}
                    .hint {{ color: #64748b; font-size: 13px; }}
                    .actions {{ margin-top: 1rem; display: flex; gap: 8px; justify-content: center; flex-wrap: wrap; }}
                    .btn {{ border: 1px solid #cbd5e1; border-radius: 6px; padding: 8px 14px; font-size: 13px; cursor: pointer; background: #f8fafc; }}
                    .btn.primary {{ background: #2563eb; color: #fff; border-color: #2563eb; }}
                </style>
            </head>
            <body>
                <div class="card">
                    <h2>{icon} {_html.escape(title)}</h2>
                    <p class="msg">{_html.escape(message)}</p>
                    <p class="hint">正在返回系统...</p>
                    <div class="actions">
                        <button class="btn primary" onclick="window.location.href='/'">返回系统首页</button>
                        <button class="btn" onclick="window.close()">关闭当前页面</button>
                    </div>
                </div>
                <script>{script}</script>
            </body>
            </html>
            """
        )

    if error:
        msg = f"飞书返回错误: {error}"
        script = f"""
            if (window.opener && !window.opener.closed) {{
                window.opener.postMessage({{ type: 'oauth_error', error: {json.dumps(msg)} }}, '*');
                setTimeout(() => window.close(), 1800);
            }} else {{
                setTimeout(() => window.location.href = '/', 1800);
            }}
        """
        return _render_page("授权失败", msg, False, script)

    if not code:
        msg = "未获取到授权码，请重新授权。"
        script = """
            if (window.opener && !window.opener.closed) {
                window.opener.postMessage({ type: 'oauth_error', error: '未获取到授权码' }, '*');
                setTimeout(() => window.close(), 1800);
            } else {
                setTimeout(() => window.location.href = '/', 1800);
            }
        """
        return _render_page("授权失败", msg, False, script)

    # 关键：在回调页直接交换 token，避免桌面模式下没有 opener 导致无法回传
    try:
        redirect_uri = "http://127.0.0.1:8000/auth/callback"
        await feishu_oauth.exchange_code_for_token(code, redirect_uri)
        script = f"""
            if (window.opener && !window.opener.closed) {{
                window.opener.postMessage({{ type: 'oauth_success', state: {json.dumps(state or "")} }}, '*');
                setTimeout(() => window.close(), 1200);
            }} else {{
                setTimeout(() => window.location.href = '/', 1200);
            }}
        """
        return _render_page("授权成功", "飞书账号授权成功。", True, script)
    except Exception as e:
        msg = f"授权码处理失败: {e}"
        script = f"""
            if (window.opener && !window.opener.closed) {{
                window.opener.postMessage({{ type: 'oauth_error', error: {json.dumps(msg)} }}, '*');
                setTimeout(() => window.close(), 2000);
            }} else {{
                setTimeout(() => window.location.href = '/', 2000);
            }}
        """
        return _render_page("授权失败", msg, False, script)


# ============ 后台任务 ============

def _read_file_preview(file_path: str, file_name: str) -> str:
    """阻塞型文件预览读取（在线程中执行）"""
    ext = Path(file_path).suffix.lower()
    try:
        if ext in [".txt", ".md", ".markdown", ".mark", ".html"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:5000]
        if ext in [".docx", ".doc"]:
            try:
                from docx import Document

                doc = Document(file_path)
                return "\n".join([p.text for p in doc.paragraphs if p.text])[:5000]
            except Exception:
                return f"Word文档: {file_name}"
        if ext in [".xlsx", ".xls", ".csv"]:
            return f"电子表格: {file_name}"
        if ext in [".pptx", ".ppt"]:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    texts.append(t)
                return "\n".join(texts)[:5000] or f"PowerPoint: {file_name}"
            except Exception:
                return f"PowerPoint: {file_name}"
        if ext == ".xmind":
            return f"XMind思维导图: {file_name}"
        if ext == ".mm":
            return f"FreeMind思维导图: {file_name}"
        if ext == ".opml":
            try:
                from xml.etree import ElementTree as ET
                tree = ET.parse(file_path)
                texts = [el.get("text", "") for el in tree.iter("outline") if el.get("text")]
                return "\n".join(texts)[:5000] or f"OPML: {file_name}"
            except Exception:
                return f"OPML: {file_name}"
        return f"文件: {file_name}"
    except Exception:
        return f"文件: {file_name}"


def _snapshot_file_status(file_status: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    return [dict(item) for item in file_status]


async def run_migration_task(task_id: str):
    """运行迁移任务（动态并发 + 异步提速 + 线程安全）"""
    from core.llm.processor import LLMProcessor
    from core.feishu.drive_api import drive_api
    from core.feishu.wiki_api import wiki_api
    from models.template import template_manager

    task = app_state["tasks"][task_id]
    files = task.get("files", [])
    duplicate_names = [
        (item.get("name") if isinstance(item, dict) else str(item))
        for item in task.get("duplicates", [])
        if (item.get("name") if isinstance(item, dict) else str(item))
    ]
    duplicate_count = len(duplicate_names)
    total_files = len(files) + duplicate_count

    if total_files == 0:
        task["status"] = "completed"
        task["progress"] = 100
        task["updated_at"] = datetime.now().isoformat()
        task["results"] = {"processed": 0, "failed": 0, "duplicate": 0, "total": 0}
        await manager.send_progress(task_id, 100, "完成（无文件）")
        return

    FEISHU_MAX_CONCURRENT_UPLOADS = 5

    processable_files_count = len(files)
    workers, worker_detail = calculate_dynamic_workers(max(processable_files_count, 1))
    workers = min(workers, FEISHU_MAX_CONCURRENT_UPLOADS)
    semaphore = asyncio.Semaphore(workers)
    runtime_started = time.perf_counter()
    chart_points: List[Dict[str, Any]] = []

    runtime = {
        "completed": duplicate_count,
        "processed": 0,
        "failed": 0,
        "duplicate": duplicate_count,
    }
    file_status_list = [
        {"name": name, "status": "duplicate", "progress": 100}
        for name in duplicate_names
    ] + [{"name": f["name"], "status": "pending", "progress": 0} for f in files]
    duplicate_offset = duplicate_count

    await manager.send_log(
        "Task",
        (
            f"开始处理 {total_files} 个文件（待处理 {processable_files_count}，重复 {duplicate_count}），"
            f"动态并发={workers} (CPU={int(worker_detail['cpu_count'])}, "
            f"可用内存≈{worker_detail['available_memory_gb']}GB)"
        ),
        "info",
    )
    app_state["system_status"] = "running"
    task["runtime"] = {"workers": workers, "detail": worker_detail}

    async def publish_runtime(phase: str):
        async with task_state_lock:
            progress = int(runtime["completed"] * 100 / total_files)
            task["progress"] = progress

            elapsed = max(0.001, time.perf_counter() - runtime_started)
            throughput = round(runtime["completed"] / elapsed, 2)
            point = {
                "x": runtime["completed"],
                "y": throughput,
                "label": datetime.now().strftime("%H:%M:%S"),
            }
            if not chart_points or chart_points[-1]["x"] != point["x"]:
                chart_points.append(point)
                if len(chart_points) > 60:
                    del chart_points[:-60]

            stats_snapshot = dict(app_state["stats"])
            status_snapshot = _snapshot_file_status(file_status_list)
            chart_snapshot = list(chart_points)

        await manager.send_progress(task_id, progress, phase)
        await manager.send_panorama(
            task_id=task_id,
            space_structure=template.get("structure", []) if template else [],
            file_status=status_snapshot,
        )
        await manager.send_chart_data(task_id=task_id, chart_type="line", data=chart_snapshot)
        await manager.send_stats(stats_snapshot)

    try:
        template_id = task.get("template_id")
        template = await template_manager.get_template(template_id) if template_id else None

        # 重复文件计入统计（避免同一任务重复启动时重复累计）
        already_counted = int(task.get("_duplicate_stats_counted", 0))
        new_duplicate_count = max(0, duplicate_count - already_counted)
        if new_duplicate_count > 0:
            async with task_state_lock:
                app_state["stats"]["duplicate"] += new_duplicate_count
            task["_duplicate_stats_counted"] = already_counted + new_duplicate_count

        for dup_name in duplicate_names:
            await manager.send_log("Upload", f"重复文件跳过: {dup_name}", "warn")

        if processable_files_count == 0:
            task["status"] = "completed"
            task["progress"] = 100
            task["updated_at"] = datetime.now().isoformat()
            task["results"] = {
                "processed": 0,
                "failed": 0,
                "duplicate": duplicate_count,
                "total": total_files,
                "workers": 0,
            }
            await manager.send_panorama(
                task_id=task_id,
                space_structure=template.get("structure", []) if template else [],
                file_status=_snapshot_file_status(file_status_list),
            )
            await manager.send_stats(dict(app_state["stats"]))
            await manager.send_progress(task_id, 100, "完成（全部重复）")
            await manager.send_log("Task", f"任务完成：全部为重复文件，共 {duplicate_count} 个", "success")
            return

        space_name = template["name"] if template else task["name"]
        space_description = template.get("description", "") if template else "由知识库迁移系统创建"

        if task.get("target_space_id"):
            space_id = task["target_space_id"]
            await manager.send_log("Feishu", f"使用已有知识空间: {space_id}", "info")
        else:
            try:
                space_id = await wiki_api.create_space(name=space_name, description=space_description)
                task["target_space_id"] = space_id
                await manager.send_log("Feishu", f"创建知识空间: {space_name} ({space_id})", "success")
            except Exception as e:
                logger.error(f"创建知识空间失败: {e}")
                await manager.send_log("Feishu", f"创建知识空间失败: {e}", "error")
                raise Exception(f"无法创建知识空间，请检查应用权限: {e}")

        node_list: List[Dict[str, Any]] = []
        node_token_map: Dict[str, str] = {}
        if template and template.get("structure"):
            await manager.send_log("Feishu", "正在根据模板结构创建文档层级...", "info")
            try:
                node_token_map = await wiki_api.create_structure(
                    space_id=space_id,
                    structure=template["structure"],
                )
                node_list = await wiki_api.get_space_nodes_flat(space_id)
                await manager.send_log(
                    "Feishu",
                    f"模板结构创建完成，节点={len(node_token_map)}，可挂载点={len(node_list)}",
                    "success",
                )
            except Exception as e:
                logger.error(f"创建模板结构失败: {e}")
                await manager.send_log("Feishu", f"创建模板结构失败: {e}", "error")

        await manager.send_panorama(
            task_id=task_id,
            space_structure=template.get("structure", []) if template else [],
            file_status=_snapshot_file_status(file_status_list),
        )

        async def process_single_file(index: int, file_info: Dict[str, Any]):
            file_failed = False
            file_processed = False
            file_partial = False
            file_name = file_info["name"]
            file_path = file_info["path"]
            status_index = duplicate_offset + index

            async with semaphore:
                async with task_state_lock:
                    file_status_list[status_index]["status"] = "processing"
                    file_status_list[status_index]["progress"] = 15
                try:
                    await publish_runtime(f"处理中 {status_index + 1}/{total_files}: {file_name}")
                except Exception:
                    pass

                try:
                    await manager.send_log("Import", f"开始处理: {file_name}", "info")

                    file_content = await asyncio.to_thread(_read_file_preview, file_path, file_name)
                    processor = LLMProcessor()

                    await manager.send_log("LLM", f"快速总结: {file_name}", "info")
                    summary_result = await processor.quick_summarize_file(file_content, file_name)
                    summary_tokens = summary_result.get("tokens", 0) if summary_result.get("success") else 0
                    file_summary = summary_result.get("summary", "")

                    async with task_state_lock:
                        if summary_tokens > 0:
                            app_state["stats"]["tokens"] += summary_tokens
                        file_status_list[status_index]["progress"] = 35

                    await manager.send_log("Import", f"导入文件: {file_name}", "info")
                    import_result = await drive_api.import_file(file_path)
                    doc_token = import_result["token"]
                    doc_type = import_result["type"]
                    doc_title = import_result["title"]

                    async with task_state_lock:
                        file_status_list[status_index]["progress"] = 65

                    parent_wiki_token = ""
                    if node_list:
                        decision_result = await processor.decide_mount_node(
                            file_summary=file_summary,
                            file_name=file_name,
                            template_structure=node_list,
                        )
                        decision_tokens = decision_result.get("tokens", 0) if decision_result.get("success") else 0
                        chosen_path = decision_result.get("node_path", "")

                        async with task_state_lock:
                            if decision_tokens > 0:
                                app_state["stats"]["tokens"] += decision_tokens

                        if chosen_path and chosen_path in node_token_map:
                            parent_wiki_token = node_token_map[chosen_path]
                        elif chosen_path and chosen_path != "根目录":
                            for path, token in node_token_map.items():
                                if chosen_path in path or path in chosen_path:
                                    parent_wiki_token = token
                                    break

                    await manager.send_log(
                        "Feishu",
                        f"移动至知识空间{'(根节点)' if not parent_wiki_token else '(指定节点)'}: {doc_title}",
                        "info",
                    )
                    try:
                        move_result = await wiki_api.move_docs_to_wiki(
                            space_id=space_id,
                            parent_wiki_token=parent_wiki_token,
                            obj_type=doc_type,
                            obj_token=doc_token,
                            title=doc_title,
                        )
                        await manager.send_log("Feishu", f"移动成功: {move_result.get('url', '')}", "success")
                        file_processed = True
                        async with task_state_lock:
                            file_status_list[status_index]["status"] = "success"
                            file_status_list[status_index]["progress"] = 100
                    except Exception as move_error:
                        logger.error(f"移动文档失败: {move_error}")
                        await manager.send_log("Feishu", f"移动失败: {move_error}", "error")
                        file_partial = True
                        async with task_state_lock:
                            file_status_list[status_index]["status"] = "partial"
                            file_status_list[status_index]["progress"] = 80

                    async with task_state_lock:
                        app_state["stats"]["api_calls"] += 4

                except Exception as file_error:
                    file_failed = True
                    logger.error(f"处理文件失败 {file_name}: {file_error}")
                    try:
                        await manager.send_log("Task", f"{file_name} 处理失败: {file_error}", "error")
                    except Exception:
                        pass
                    async with task_state_lock:
                        file_status_list[status_index]["status"] = "failed"
                        file_status_list[status_index]["progress"] = 0
                finally:
                    async with task_state_lock:
                        runtime["completed"] += 1
                        if file_processed:
                            runtime["processed"] += 1
                            app_state["stats"]["processed"] += 1
                        if file_failed or file_partial:
                            runtime["failed"] += 1
                            app_state["stats"]["failed"] += 1

                    try:
                        await publish_runtime(f"已完成 {runtime['completed']}/{total_files}")
                    except Exception:
                        pass

        await asyncio.gather(
            *(process_single_file(idx, info) for idx, info in enumerate(files)),
            return_exceptions=True,
        )

        # Bug-fix: send definitive final panorama + stats after all files complete
        await manager.send_panorama(
            task_id=task_id,
            space_structure=template.get("structure", []) if template else [],
            file_status=_snapshot_file_status(file_status_list),
        )
        await manager.send_stats(dict(app_state["stats"]))

        task["status"] = "completed"
        task["progress"] = 100
        task["updated_at"] = datetime.now().isoformat()
        task["results"] = {
            "processed": runtime["processed"],
            "failed": runtime["failed"],
            "duplicate": runtime["duplicate"],
            "total": total_files,
            "workers": workers,
        }
        await manager.send_progress(task_id, 100, "完成")
        await manager.send_log(
            "Task",
            (
                f"任务完成: 成功 {runtime['processed']}, 重复 {runtime['duplicate']}, "
                f"失败 {runtime['failed']}, 并发 {workers}"
            ),
            "success",
        )

    except Exception as e:
        logger.error(f"任务执行失败: {e}")
        task["status"] = "error"
        task["error"] = str(e)
        await manager.send_log("Task", f"任务失败: {str(e)}", "error")
    finally:
        app_state["system_status"] = "ready"


# ============ 静态文件 ============

# 挂载前端静态文件
frontend_dir = BASE_DIR / "frontend"
if frontend_dir.exists():
    app.mount("/static", StaticFiles(directory=str(frontend_dir)), name="static")


if __name__ == "__main__":
    import uvicorn
    config = get_config()
    uvicorn.run(
        "main:app",
        host=config.host,
        port=config.port,
        reload=config.debug
    )
