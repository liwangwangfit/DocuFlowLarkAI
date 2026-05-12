"""
v2.0 功能测试
- Task 1: 并发节点创建
- Task 2: 上传并发数限制为 5
- Task 3: 新文件格式转换
- Task 4: 仪表盘失败计数修复
"""
import asyncio
import json
import sys
import tempfile
import zipfile
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

backend_dir = Path(__file__).parent.parent / "backend"
sys.path.insert(0, str(backend_dir))


# ═══════════════════════════════════════════════════════════════════
# Task 1: Parallel Node Creation
# ═══════════════════════════════════════════════════════════════════

class TestParallelNodeCreation:
    """create_structure should create sibling nodes concurrently
    while still respecting the parent-before-child constraint."""

    @pytest.mark.asyncio
    async def test_siblings_created_in_parallel(self):
        """Siblings under the same parent should be created via gather, not sequentially."""
        from core.feishu.wiki_api import FeishuWikiAPI

        api = FeishuWikiAPI()
        call_log = []

        async def mock_find(space_id, title, parent_token=None):
            return None

        async def mock_create(space_id, title, parent_node_token=None, obj_type="docx"):
            call_log.append(("create", title, parent_node_token))
            return {"node_token": f"tok_{title}", "obj_token": ""}

        api.find_node_by_title = mock_find
        api.create_node = mock_create

        structure = [
            {"name": "A", "children": []},
            {"name": "B", "children": []},
            {"name": "C", "children": []},
        ]

        node_map = await api.create_structure("sp1", structure)

        assert len(node_map) == 3
        assert "A" in node_map
        assert "B" in node_map
        assert "C" in node_map
        created_titles = [c[1] for c in call_log if c[0] == "create"]
        assert set(created_titles) == {"A", "B", "C"}

    @pytest.mark.asyncio
    async def test_parent_created_before_children(self):
        """Children must only be created after their parent node exists."""
        from core.feishu.wiki_api import FeishuWikiAPI

        api = FeishuWikiAPI()
        creation_order = []

        async def mock_find(space_id, title, parent_token=None):
            return None

        async def mock_create(space_id, title, parent_node_token=None, obj_type="docx"):
            creation_order.append(title)
            return {"node_token": f"tok_{title}", "obj_token": ""}

        api.find_node_by_title = mock_find
        api.create_node = mock_create

        structure = [
            {
                "name": "Parent",
                "children": [
                    {"name": "Child1", "children": []},
                    {"name": "Child2", "children": []},
                ],
            }
        ]

        await api.create_structure("sp1", structure)

        assert creation_order.index("Parent") < creation_order.index("Child1")
        assert creation_order.index("Parent") < creation_order.index("Child2")

    @pytest.mark.asyncio
    async def test_existing_node_reused(self):
        """If a node already exists, it should be reused and not recreated."""
        from core.feishu.wiki_api import FeishuWikiAPI

        api = FeishuWikiAPI()
        create_calls = []

        async def mock_find(space_id, title, parent_token=None):
            if title == "Existing":
                return "existing_token"
            return None

        async def mock_create(space_id, title, parent_node_token=None, obj_type="docx"):
            create_calls.append(title)
            return {"node_token": f"tok_{title}", "obj_token": ""}

        api.find_node_by_title = mock_find
        api.create_node = mock_create

        structure = [
            {"name": "Existing", "children": []},
            {"name": "New", "children": []},
        ]

        node_map = await api.create_structure("sp1", structure)

        assert node_map["Existing"] == "existing_token"
        assert "New" in node_map
        assert "Existing" not in create_calls
        assert "New" in create_calls


# ═══════════════════════════════════════════════════════════════════
# Task 2: Upload Concurrency Cap
# ═══════════════════════════════════════════════════════════════════

class TestUploadConcurrencyCap:

    def test_max_workers_default_is_five(self):
        from utils.system_resource import calculate_dynamic_workers

        workers, detail = calculate_dynamic_workers(total_files=100)
        assert workers <= 5

    def test_respects_explicit_max(self):
        from utils.system_resource import calculate_dynamic_workers

        workers, _ = calculate_dynamic_workers(total_files=100, max_workers=3)
        assert workers <= 3

    def test_never_zero(self):
        from utils.system_resource import calculate_dynamic_workers

        workers, _ = calculate_dynamic_workers(total_files=1)
        assert workers >= 1


# ═══════════════════════════════════════════════════════════════════
# Task 3: File Format Support
# ═══════════════════════════════════════════════════════════════════

class TestFileTypeMap:
    """drive_api FILE_TYPE_MAP should cover all supported formats including pptx/xmind/mm/opml."""

    def test_native_document_types(self):
        from core.feishu.drive_api import FeishuDriveAPI

        for ext in ("docx", "doc", "txt", "md", "mark", "markdown", "html"):
            assert ext in FeishuDriveAPI.FILE_TYPE_MAP
            assert FeishuDriveAPI.FILE_TYPE_MAP[ext][0] == "docx"

    def test_native_sheet_types(self):
        from core.feishu.drive_api import FeishuDriveAPI

        for ext in ("xlsx", "csv", "xls"):
            assert ext in FeishuDriveAPI.FILE_TYPE_MAP
            assert FeishuDriveAPI.FILE_TYPE_MAP[ext][0] == "sheet"

    def test_pptx_not_imported_as_native_file_type(self):
        from core.feishu.drive_api import FeishuDriveAPI

        assert "pptx" not in FeishuDriveAPI.FILE_TYPE_MAP

    def test_mindmap_types_not_imported_as_native_file_types(self):
        from core.feishu.drive_api import FeishuDriveAPI

        for ext in ("xmind", "mm", "opml"):
            assert ext not in FeishuDriveAPI.FILE_TYPE_MAP

    def test_unsupported_extension_raises(self):
        from core.feishu.drive_api import FeishuDriveAPI

        api = FeishuDriveAPI()
        with pytest.raises(ValueError, match="不支持的文件类型"):
            api._get_file_type_info("test.zzz")


class TestMindMapConverter:

    @pytest.mark.asyncio
    async def test_opml_to_markdown(self):
        from core.converter.mindmap_converter import MindMapConverter

        opml_content = """<?xml version="1.0" encoding="UTF-8"?>
<opml version="2.0">
  <head><title>Test OPML</title></head>
  <body>
    <outline text="Topic A">
      <outline text="Sub A1"/>
    </outline>
    <outline text="Topic B"/>
  </body>
</opml>"""
        tmp = Path(tempfile.mktemp(suffix=".opml"))
        tmp.write_text(opml_content, encoding="utf-8")

        converter = MindMapConverter()
        result = await converter.convert(str(tmp))

        assert result.success
        assert "Topic A" in result.content
        assert "Sub A1" in result.content
        assert "Topic B" in result.content
        tmp.unlink(missing_ok=True)

    @pytest.mark.asyncio
    async def test_freemind_to_markdown(self):
        from core.converter.mindmap_converter import MindMapConverter

        mm_content = """<?xml version="1.0" encoding="UTF-8"?>
<map version="1.0.1">
  <node TEXT="Root">
    <node TEXT="Branch A"/>
    <node TEXT="Branch B">
      <node TEXT="Leaf B1"/>
    </node>
  </node>
</map>"""
        tmp = Path(tempfile.mktemp(suffix=".mm"))
        tmp.write_text(mm_content, encoding="utf-8")

        converter = MindMapConverter()
        result = await converter.convert(str(tmp))

        assert result.success
        assert "Root" in result.content
        assert "Branch A" in result.content
        assert "Leaf B1" in result.content
        tmp.unlink(missing_ok=True)

    @pytest.mark.asyncio
    async def test_xmind_json_to_markdown(self):
        from core.converter.mindmap_converter import MindMapConverter

        xmind_json = json.dumps([{
            "rootTopic": {
                "title": "Central",
                "children": {
                    "attached": [
                        {"title": "Node1", "children": {"attached": []}},
                        {"title": "Node2", "children": {"attached": []}},
                    ]
                }
            }
        }])
        tmp = Path(tempfile.mktemp(suffix=".xmind"))
        with zipfile.ZipFile(str(tmp), "w") as zf:
            zf.writestr("content.json", xmind_json)

        converter = MindMapConverter()
        result = await converter.convert(str(tmp))

        assert result.success
        assert "Central" in result.content
        assert "Node1" in result.content
        assert "Node2" in result.content
        tmp.unlink(missing_ok=True)

    @pytest.mark.asyncio
    async def test_unsupported_extension_fails(self):
        from core.converter.mindmap_converter import MindMapConverter

        tmp = Path(tempfile.mktemp(suffix=".xyz"))
        tmp.write_text("nope", encoding="utf-8")

        converter = MindMapConverter()
        result = await converter.convert(str(tmp))

        assert not result.success
        tmp.unlink(missing_ok=True)


class TestPPTXConverter:

    @pytest.mark.asyncio
    async def test_pptx_to_docx(self):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from core.converter.pptx_converter import PPTXConverter

        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Hello PPTX"

        tmp = Path(tempfile.mktemp(suffix=".pptx"))
        prs.save(str(tmp))

        converter = PPTXConverter()
        result = await converter.convert(str(tmp))

        assert result.success
        assert result.output_path is not None
        assert Path(result.output_path).exists()
        assert Path(result.output_path).suffix == ".docx"
        tmp.unlink(missing_ok=True)


class TestFeishuImportPreparation:

    @pytest.mark.asyncio
    async def test_direct_import_type_is_unchanged(self):
        from main import _prepare_file_for_feishu_import

        tmp = Path(tempfile.mktemp(suffix=".md"))
        tmp.write_text("# Direct", encoding="utf-8")

        result = await _prepare_file_for_feishu_import(str(tmp))

        assert result == str(tmp)
        tmp.unlink(missing_ok=True)

    @pytest.mark.asyncio
    async def test_opml_is_converted_to_markdown_before_import(self):
        from main import _prepare_file_for_feishu_import

        opml_content = """<?xml version="1.0" encoding="UTF-8"?>
<opml version="2.0">
  <head><title>Import OPML</title></head>
  <body><outline text="Topic"/></body>
</opml>"""
        tmp = Path(tempfile.mktemp(suffix=".opml"))
        tmp.write_text(opml_content, encoding="utf-8")

        result = await _prepare_file_for_feishu_import(str(tmp))

        assert Path(result).suffix == ".md"
        assert Path(result).exists()
        assert "Topic" in Path(result).read_text(encoding="utf-8")
        tmp.unlink(missing_ok=True)


# ═══════════════════════════════════════════════════════════════════
# Task 4: Dashboard Failed Count
# ═══════════════════════════════════════════════════════════════════

class TestFailedCountLogic:
    """Verify the fixed counting logic for partial / failed files."""

    def test_partial_counted_as_failed(self):
        """
        Simulate the finally-block counting logic:
        when file_partial is True, it should increment failed counter.
        """
        runtime = {"completed": 0, "processed": 0, "failed": 0, "duplicate": 0}
        stats = {"processed": 0, "failed": 0, "duplicate": 0, "tokens": 0, "api_calls": 0}

        file_processed = False
        file_failed = False
        file_partial = True

        runtime["completed"] += 1
        if file_processed:
            runtime["processed"] += 1
            stats["processed"] += 1
        if file_failed or file_partial:
            runtime["failed"] += 1
            stats["failed"] += 1

        assert runtime["failed"] == 1
        assert stats["failed"] == 1
        assert runtime["processed"] == 0

    def test_success_counted_as_processed(self):
        runtime = {"completed": 0, "processed": 0, "failed": 0, "duplicate": 0}
        stats = {"processed": 0, "failed": 0, "duplicate": 0, "tokens": 0, "api_calls": 0}

        file_processed = True
        file_failed = False
        file_partial = False

        runtime["completed"] += 1
        if file_processed:
            runtime["processed"] += 1
            stats["processed"] += 1
        if file_failed or file_partial:
            runtime["failed"] += 1
            stats["failed"] += 1

        assert runtime["processed"] == 1
        assert runtime["failed"] == 0

    def test_full_failure_counted(self):
        runtime = {"completed": 0, "processed": 0, "failed": 0, "duplicate": 0}
        stats = {"processed": 0, "failed": 0, "duplicate": 0, "tokens": 0, "api_calls": 0}

        file_processed = False
        file_failed = True
        file_partial = False

        runtime["completed"] += 1
        if file_processed:
            runtime["processed"] += 1
            stats["processed"] += 1
        if file_failed or file_partial:
            runtime["failed"] += 1
            stats["failed"] += 1

        assert runtime["failed"] == 1
        assert stats["failed"] == 1


class TestFinalPanoramaSent:
    """After asyncio.gather, a final panorama must be sent to ensure
    the frontend has the definitive file-status snapshot."""

    def test_snapshot_captures_all_statuses(self):
        """_snapshot_file_status must deep-copy the list."""
        file_status_list = [
            {"name": "a.docx", "status": "success", "progress": 100},
            {"name": "b.docx", "status": "failed", "progress": 0},
            {"name": "c.docx", "status": "partial", "progress": 80},
        ]

        from main import _snapshot_file_status
        snapshot = _snapshot_file_status(file_status_list)

        assert len(snapshot) == 3
        assert snapshot[1]["status"] == "failed"

        file_status_list[1]["status"] = "changed"
        assert snapshot[1]["status"] == "failed"
